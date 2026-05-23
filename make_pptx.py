from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colors ────────────────────────────────────────────────
NAVY       = RGBColor(0x09, 0x1A, 0x3A)   # deep navy
BLUE       = RGBColor(0x1A, 0x56, 0xC8)   # primary blue
ACCENT     = RGBColor(0x1F, 0xB7, 0xFF)   # light cyan
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTED      = RGBColor(0xB0, 0xC4, 0xDE)
GREEN      = RGBColor(0x1A, 0xC2, 0x72)
GOLD       = RGBColor(0xF7, 0xB7, 0x31)
DARK_PANEL = RGBColor(0x0D, 0x24, 0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

LOGO = r"assets\images\mbongo_logo.png"
LOGO_EXISTS = os.path.exists(LOGO)

SCREENS = {
    "login":   r"emulator-mbongo-after-login.png",
    "home":    r"emulator-mbongo-screen.png",
    "kyc":     r"emulator-mbongo-kyc-fixed.png",
    "kyc2":    r"emulator-mbongo-kyc-fixed-2.png",
}
SCREENS = {k: v for k, v in SCREENS.items() if os.path.exists(v)}


def bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def panel(slide, l, t, w, h, color=DARK_PANEL, radius=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(0.5)
    return shape


def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def accent_bar(slide, l, t, w=0.05, h=0.5):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def chip(slide, label, l, t, color=BLUE):
    w, h = 1.6, 0.36
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE


# ══════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)

# big gradient bar on left
bar = sl.shapes.add_shape(1, 0, 0, Inches(0.12), prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

# Logo
if LOGO_EXISTS:
    sl.shapes.add_picture(LOGO, Inches(1.0), Inches(1.6), height=Inches(1.2))
    logo_offset = 1.6 + 1.4
else:
    txt(sl, "MBONGO", 1.0, 1.6, 5, 1.4, 56, True, ACCENT, PP_ALIGN.LEFT)
    logo_offset = 1.6 + 1.4

txt(sl, "Application Financière Mobile", 1.0, logo_offset, 9, 0.8,
    24, False, MUTED)
txt(sl, "Votre portefeuille, partout, tout le temps.", 1.0, logo_offset + 0.7,
    10, 0.7, 18, False, WHITE)

# Right side — feature chips
chips = [
    ("Portefeuille CDF & USD", 8.5, 1.8, BLUE),
    ("Recharge mobile",        8.5, 2.3, BLUE),
    ("Transferts & Retraits",  8.5, 2.8, BLUE),
    ("Cartes virtuelles",      8.5, 3.3, BLUE),
    ("Réabonnement TV",        8.5, 3.8, BLUE),
    ("Vérification KYC",       8.5, 4.3, BLUE),
    ("Backoffice Admin",       8.5, 4.8, RGBColor(0x0D, 0x46, 0x8A)),
]
for label, l, t, color in chips:
    chip(sl, label, l, t, color)

# Bottom tag
txt(sl, "© 2026 MBONGO — Plateforme Fintech Afrique Centrale",
    1.0, 6.9, 12, 0.4, 9, False, MUTED, PP_ALIGN.LEFT)

# Slide number dot
dot = sl.shapes.add_shape(9, Inches(12.8), Inches(7.1), Inches(0.3), Inches(0.3))
dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
dot.line.fill.background()


# ══════════════════════════════════════════════════════════
# SLIDE 2 — Vue d'ensemble
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
accent_bar(sl, 0, 0)

txt(sl, "Vue d'ensemble", 0.4, 0.25, 10, 0.6, 11, True, ACCENT)
txt(sl, "Qu'est-ce que MBONGO ?", 0.4, 0.9, 12, 0.7, 26, True, WHITE)

body = (
    "MBONGO est une application fintech mobile conçue pour l'Afrique Centrale "
    "(RDC & diaspora). Elle offre à ses utilisateurs un portefeuille électronique "
    "multi-devises (CDF / USD), des services de transfert d'argent, de recharge "
    "mobile, de paiement de factures et de cartes virtuelles Visa/Mastercard — "
    "le tout depuis leur smartphone."
)
txt(sl, body, 0.4, 1.75, 12.5, 1.3, 13, False, MUTED)

# 4 pillars
pillars = [
    ("💳", "Portefeuille", "Solde multi-devises, historique complet, carte numérique"),
    ("📲", "Services",     "Recharge, TV, transfert international, paiement factures"),
    ("🔒", "Sécurité",     "KYC obligatoire, OTP, PIN chiffré, 2FA admin"),
    ("🖥️", "Backoffice",   "Admin Next.js, 59 modules, API NestJS sur Vercel"),
]
for i, (icon, title, desc) in enumerate(pillars):
    l = 0.3 + i * 3.2
    p = panel(sl, l, 3.3, 3.0, 2.8)
    txt(sl, icon,  l+0.15, 3.4,  0.7, 0.7, 26)
    txt(sl, title, l+0.15, 4.1,  2.7, 0.5, 14, True,  WHITE)
    txt(sl, desc,  l+0.15, 4.65, 2.7, 1.2, 11, False, MUTED)


# ══════════════════════════════════════════════════════════
# SLIDE 3 — Stack technique
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
accent_bar(sl, 0, 0)

txt(sl, "Architecture technique", 0.4, 0.25, 10, 0.6, 11, True, ACCENT)
txt(sl, "Stack & Infrastructure", 0.4, 0.9, 12, 0.7, 26, True, WHITE)

layers = [
    ("Application Mobile",  "Flutter 3 · Dart · Riverpod · Dio · FCM Push",      BLUE,  1.9),
    ("Backend API",         "NestJS 11 · Prisma 6 · PostgreSQL · JWT · Vercel",   GOLD,  2.95),
    ("Admin Backoffice",    "Next.js 16 · Tailwind CSS · 59 modules · Vercel",    GREEN, 4.0),
    ("Infrastructure",      "Vercel Serverless · GitHub Actions · CI/CD APK",     ACCENT,5.05),
]
for label, detail, color, t in layers:
    bar2 = sl.shapes.add_shape(1, Inches(0.4), Inches(t), Inches(0.08), Inches(0.75))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = color
    bar2.line.fill.background()
    p2 = panel(sl, 0.6, t, 11.9, 0.75, RGBColor(0x0D, 0x24, 0x4A))
    txt(sl, label,  0.8, t+0.04, 3.5, 0.4, 13, True,  color)
    txt(sl, detail, 4.5, t+0.04, 8.0, 0.4, 12, False, MUTED)

txt(sl, "Toutes les API sont sécurisées par Bearer Token (JWT). Le backend tourne "
        "en serverless sur Vercel avec Prisma connecté à PostgreSQL.",
    0.4, 6.2, 12.5, 0.8, 11, False, MUTED)


# ══════════════════════════════════════════════════════════
# SLIDE 4 — Fonctionnalités mobiles
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
accent_bar(sl, 0, 0)

txt(sl, "Application Mobile", 0.4, 0.25, 10, 0.6, 11, True, ACCENT)
txt(sl, "Fonctionnalités utilisateur", 0.4, 0.9, 12, 0.7, 26, True, WHITE)

features = [
    ("Portefeuille",        ["Solde CDF & USD en temps réel", "Historique transactions", "Masquage solde"]),
    ("Transferts",          ["Envoi entre comptes MBONGO", "Demande d'argent (QR)", "International (remittance)"]),
    ("Recharge & Factures", ["Airtime Vodacom / Airtel / Orange / Africell", "Réabonnement TV (Canal+, StarTimes…)", "Paiement de factures"]),
    ("Cartes Virtuelles",   ["Création carte Visa/Mastercard", "Rechargement, suspension", "Historique transactions carte"]),
    ("Sécurité & KYC",      ["Vérification identité (recto/verso)", "PIN chiffré + OTP", "Garde KYC sur opérations sensibles"]),
    ("Retrait & Dépôt",     ["Mobile Money (M-Pesa, Orange, Airtel)", "Guichet agent + DAB", "Dépôt via méthodes configurables"]),
]
for i, (title, bullets) in enumerate(features):
    col = i % 3
    row = i // 3
    l = 0.3 + col * 4.35
    t = 2.0 + row * 2.4
    p2 = panel(sl, l, t, 4.1, 2.15)
    txt(sl, title, l+0.15, t+0.1, 3.8, 0.45, 13, True, ACCENT)
    for j, b in enumerate(bullets):
        txt(sl, f"• {b}", l+0.15, t+0.55+j*0.48, 3.8, 0.45, 11, False, MUTED)


# ══════════════════════════════════════════════════════════
# SLIDE 5 — Captures d'écran mobile
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
accent_bar(sl, 0, 0)

txt(sl, "Interface Mobile", 0.4, 0.25, 10, 0.6, 11, True, ACCENT)
txt(sl, "Aperçu de l'application", 0.4, 0.9, 12, 0.7, 26, True, WHITE)

screen_positions = [
    ("login", 0.4,  1.85, "Connexion"),
    ("home",  3.6,  1.85, "Accueil"),
    ("kyc",   6.8,  1.85, "KYC"),
    ("kyc2",  10.0, 1.85, "Dossier KYC"),
]

placed = 0
for key, l, t, caption in screen_positions:
    if key in SCREENS:
        try:
            sl.shapes.add_picture(SCREENS[key], Inches(l), Inches(t), height=Inches(4.5))
            placed += 1
        except Exception:
            pass
    # caption
    txt(sl, caption, l, t+4.6, 2.8, 0.4, 11, True, MUTED, PP_ALIGN.CENTER)

if placed == 0:
    txt(sl, "Captures d'écran disponibles après le build APK release.",
        2.0, 3.0, 9, 1.0, 16, False, MUTED, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 6 — Backoffice Admin
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
accent_bar(sl, 0, 0)

txt(sl, "Backoffice Admin", 0.4, 0.25, 10, 0.6, 11, True, ACCENT)
txt(sl, "Panneau d'administration MBONGO", 0.4, 0.9, 12, 0.7, 26, True, WHITE)

groups = [
    ("Pilotage",              "Tableau de bord, KPIs, flux financiers temps réel"),
    ("Finances",              "Wallet, dépôts, retraits, transferts, frais, devises"),
    ("Sécurité",              "KYC, admins, rôles, audit logs, restrictions pays"),
    ("Services financiers",   "Cartes, recharge, TV, liens paiement, marchands, agents"),
    ("Paramètres plateforme", "SMTP, SEO, modules, langues, images, splash, pages"),
    ("Communication",         "Notifications push, newsletter, messages contact"),
    ("Compte admin",          "Profil, PIN, 2FA TOTP, RGPD, logs serveur & erreurs"),
]
for i, (title, desc) in enumerate(groups):
    col = i % 2
    row = i // 2
    l = 0.3 + col * 6.55
    t = 2.0 + row * 1.15
    if i == 6:
        l = 0.3 + 6.55/2
    p2 = panel(sl, l, t, 6.2, 0.95)
    txt(sl, title, l+0.15, t+0.07, 2.8, 0.38, 12, True,  WHITE)
    txt(sl, desc,  l+0.15, t+0.47, 6.0, 0.38, 10, False, MUTED)

txt(sl, "59 modules · Déployé sur Vercel · Sidebar avec recherche et groupes repliables",
    0.4, 6.8, 12.5, 0.5, 10, False, MUTED, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 7 — Sécurité & KYC
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
accent_bar(sl, 0, 0)

txt(sl, "Sécurité & Conformité", 0.4, 0.25, 10, 0.6, 11, True, ACCENT)
txt(sl, "KYC · Authentification · Protection", 0.4, 0.9, 12, 0.7, 26, True, WHITE)

securities = [
    ("🪪", "KYC Obligatoire",     "Pièce d'identité recto/verso + feuille manuscrite. Validation admin. Blocage des opérations sensibles avant validation."),
    ("🔑", "Authentification",     "OTP SMS à la connexion. PIN chiffré. Authentification 2FA TOTP pour les admins (Google Authenticator)."),
    ("🛡️", "Contrôle d'accès",    "Rôles et permissions granulaires pour les sous-admins. Audit logs complet. Restriction par pays et IP."),
    ("📊", "Monitoring",          "Logs d'erreurs, logs serveur, audit des actions admin. Export des données. Alertes configurables."),
]
for i, (icon, title, desc) in enumerate(securities):
    col = i % 2
    row = i // 2
    l = 0.3 + col * 6.55
    t = 2.1 + row * 2.2
    p2 = panel(sl, l, t, 6.2, 2.0)
    txt(sl, icon,  l+0.2, t+0.15, 0.8, 0.7, 28)
    txt(sl, title, l+1.1, t+0.15, 5.0, 0.5, 14, True, ACCENT)
    txt(sl, desc,  l+0.2, t+0.85, 5.8, 1.0, 11, False, MUTED)


# ══════════════════════════════════════════════════════════
# SLIDE 8 — Services financiers
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
accent_bar(sl, 0, 0)

txt(sl, "Services Financiers", 0.4, 0.25, 10, 0.6, 11, True, ACCENT)
txt(sl, "Offre complète de services", 0.4, 0.9, 12, 0.7, 26, True, WHITE)

services_data = [
    ("Mobile Money",        "M-Pesa · Orange Money · Airtel Money · AfriMoney", GREEN),
    ("Cartes Virtuelles",   "Émission Visa/Mastercard via API partenaire",      BLUE),
    ("Réabonnement TV",     "Canal+ · StarTimes · EasyTV · DRTV · AzamTV",     GOLD),
    ("Recharge Mobile",     "Vodacom · Airtel · Orange · Africell",             ACCENT),
    ("Transfert Intl",      "Corridors de remittance, partenaires API",         RGBColor(0xAB, 0x47, 0xBC)),
    ("Paiement Factures",   "Eau · Électricité · Internet · Assurance · École", RGBColor(0x26, 0xC6, 0xDA)),
    ("Liens de Paiement",   "Génération lien, paiement carte ou wallet",        RGBColor(0xFF, 0x70, 0x43)),
    ("Cartes Cadeaux",      "Catalogue produits, achat, livraison code",        RGBColor(0x66, 0xBB, 0x6A)),
]
for i, (title, detail, color) in enumerate(services_data):
    col = i % 4
    row = i // 4
    l = 0.25 + col * 3.28
    t = 2.1 + row * 1.9
    p2 = panel(sl, l, t, 3.1, 1.7)
    dot2 = sl.shapes.add_shape(9, Inches(l+0.15), Inches(t+0.15), Inches(0.12), Inches(0.12))
    dot2.fill.solid(); dot2.fill.fore_color.rgb = color
    dot2.line.fill.background()
    txt(sl, title,  l+0.4, t+0.1,  2.6, 0.4, 12, True,  color)
    txt(sl, detail, l+0.15, t+0.55, 2.8, 0.9, 10, False, MUTED)


# ══════════════════════════════════════════════════════════
# SLIDE 9 — Roadmap
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
accent_bar(sl, 0, 0)

txt(sl, "Roadmap", 0.4, 0.25, 10, 0.6, 11, True, ACCENT)
txt(sl, "Ce qui est fait · Ce qui vient", 0.4, 0.9, 12, 0.7, 26, True, WHITE)

done = [
    "Application Flutter complète (Android/iOS)",
    "Backend NestJS 11 + Prisma 6 (Vercel)",
    "Authentification OTP + PIN + 2FA admin",
    "KYC complet (pièces + feuille manuscrite)",
    "Portefeuille CDF/USD + transactions",
    "Retraits Mobile Money + Guichet + DAB",
    "Cartes virtuelles + recharge airtime + TV",
    "Backoffice admin 59 modules (Vercel)",
    "Push notifications FCM",
    "Transferts internationaux (remittance)",
]

coming = [
    "Cartes cadeaux (module en développement)",
    "Build iOS / App Store",
    "Intégration SMTP réelle (envoi emails)",
    "Paiements en ligne (e-commerce)",
    "Portefeuille crypto",
    "Programme de parrainage actif",
]

panel(sl, 0.3, 2.0, 6.1, 4.8, DARK_PANEL)
txt(sl, "✅  Réalisé", 0.45, 2.1, 5.8, 0.45, 13, True, GREEN)
for i, item in enumerate(done):
    txt(sl, f"  {item}", 0.45, 2.6+i*0.42, 5.8, 0.38, 10.5, False, WHITE)

panel(sl, 6.7, 2.0, 6.1, 4.8, DARK_PANEL)
txt(sl, "🔜  À venir", 6.85, 2.1, 5.8, 0.45, 13, True, GOLD)
for i, item in enumerate(coming):
    txt(sl, f"  {item}", 6.85, 2.6+i*0.42, 5.8, 0.38, 10.5, False, WHITE)


# ══════════════════════════════════════════════════════════
# SLIDE 10 — Contact / Clôture
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)

bar3 = sl.shapes.add_shape(1, 0, 0, Inches(0.12), prs.slide_height)
bar3.fill.solid(); bar3.fill.fore_color.rgb = ACCENT
bar3.line.fill.background()

if LOGO_EXISTS:
    sl.shapes.add_picture(LOGO, Inches(5.2), Inches(1.5), height=Inches(1.4))

txt(sl, "MBONGO", 3.5, 3.2, 6.5, 1.0, 44, True, WHITE, PP_ALIGN.CENTER)
txt(sl, "Votre partenaire financier digital en Afrique Centrale",
    2.0, 4.3, 9.5, 0.6, 16, False, MUTED, PP_ALIGN.CENTER)

divider = sl.shapes.add_shape(1, Inches(4.0), Inches(5.1), Inches(5.3), Pt(1))
divider.fill.solid(); divider.fill.fore_color.rgb = ACCENT
divider.line.fill.background()

txt(sl, "admin.mbongo.app  ·  mbongo-admin.vercel.app",
    2.0, 5.4, 9.5, 0.5, 13, False, ACCENT, PP_ALIGN.CENTER)
txt(sl, "Backend : mbongo-backend.vercel.app",
    2.0, 5.95, 9.5, 0.5, 11, False, MUTED, PP_ALIGN.CENTER)
txt(sl, "© 2026 MBONGO — Plateforme Fintech Afrique Centrale",
    2.0, 6.8, 9.5, 0.4, 9, False, MUTED, PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────
OUT = r"C:\Users\Silikin Store\Desktop\MBONGO_Presentation.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
